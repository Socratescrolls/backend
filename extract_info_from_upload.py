# extract_info_from_upload.py (No Changes needed)
import os
from dotenv import load_dotenv
from pptx import Presentation
from PIL import Image
import io
import base64
from openai import OpenAI
import json
import fitz  # PyMuPDF for PDF processing
import tempfile

load_dotenv()  # Load environment variables from .env file
open_api_key = os.getenv("OPENAI_API_KEY")

def get_file_type(file_path):
    """
    Determine the file type based on extension
    """
    extension = os.path.splitext(file_path)[1].lower()
    if extension == '.pdf':
        return 'pdf'
    elif extension in ['.ppt', '.pptx']:
        return 'presentation'
    else:
        raise ValueError(f"Unsupported file type: {extension}")

def extract_pdf_content(pdf_path, client):
    """
    Extract content from PDF including text and images.
    Analyzes images using OpenAI's Vision model.
    """
    doc = fitz.open(pdf_path)
    page_contents = []
    
    for page_number, page in enumerate(doc, 1):
        print(f"Processing PDF page {page_number}...")
        page_content = f"Page {page_number}:\n"
        
        # Extract text
        text = page.get_text()
        if text.strip():
            page_content += "Text content:\n" + text.strip() + "\n"
        
        # Extract images
        image_descriptions = []
        image_list = page.get_images()
        
        for img_index, img in enumerate(image_list):
            try:
                # Get image data
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                
                # Convert to PIL Image
                image = Image.open(io.BytesIO(image_bytes))
                
                # Convert image to base64
                buffered = io.BytesIO()
                image.save(buffered, format="PNG")
                img_str = base64.b64encode(buffered.getvalue()).decode()
                
                # Analyze image using OpenAI's Vision model
                response = client.chat.completions.create(
                    model="gpt-4o-mini",
                    messages=[
                        {
                            "role": "user",
                            "content": [
                                {
                                    "type": "text",
                                    "text": "Describe this image in detail, extracting the text as well as the images."
                                },
                                {
                                    "type": "image_url",
                                    "image_url": {
                                        "url": f"data:image/png;base64,{img_str}"
                                    }
                                }
                            ]
                        }
                    ],
                    max_tokens=300
                )
                
                image_description = response.choices[0].message.content
                image_descriptions.append(image_description)
                
            except Exception as e:
                print(f"Error processing image {img_index + 1} on page {page_number}: {str(e)}")
                continue
        
        if image_descriptions:
            page_content += "\nImage content:\n" + "\n".join(image_descriptions) + "\n"
        
        page_contents.append(page_content)
    
    doc.close()
    return page_contents

def extract_ppt_content(ppt_path, client):
    """
    Extract content from PowerPoint slides including text and images.
    Analyzes images using OpenAI's Vision model.
    """
    # Load presentation
    prs = Presentation(ppt_path)
    slide_contents = []
    
    # Process each slide
    for slide_number, slide in enumerate(prs.slides, 1):
        print(f"Processing slide {slide_number}...")
        slide_content = f"Slide {slide_number}:\n"
        
        # Extract text
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        
        if texts:
            slide_content += "Text content:\n" + "\n".join(texts) + "\n"
        
        # Extract and analyze images
        image_descriptions = []
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                try:
                    # Get image data
                    image_bytes = shape.image.blob
                    image = Image.open(io.BytesIO(image_bytes))
                    
                    # Convert image to base64
                    buffered = io.BytesIO()
                    image.save(buffered, format="PNG")
                    img_str = base64.b64encode(buffered.getvalue()).decode()
                    
                    # Analyze image using OpenAI's Vision model
                    response = client.chat.completions.create(
                        model="gpt-4o-mini",
                        messages=[
                            {
                                "role": "user",
                                "content": [
                                    {
                                        "type": "text",
                                        "text": "Describe this image in detail, extracting the text as well as the images."
                                    },
                                    {
                                        "type": "image_url",
                                        "image_url": {
                                            "url": f"data:image/png;base64,{img_str}"
                                        }
                                    }
                                ]
                            }
                        ],
                        max_tokens=300
                    )
                    
                    image_description = response.choices[0].message.content
                    image_descriptions.append(image_description)
                    
                except Exception as e:
                    print(f"Error processing image: {str(e)}")
                    continue
        
        if image_descriptions:
            slide_content += "\nImage content:\n" + "\n".join(image_descriptions) + "\n"
        
        slide_contents.append(slide_content)
        
    return slide_contents

def save_to_file(contents, output_path):
    """
    Save the extracted content to a text file.
    """
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write('\n\n'.join(contents))

def process_document(file_path):
    """
    Process either PDF or PowerPoint document and extract content
    """
    # Get API key from environment variable
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        raise ValueError("OpenAI API key not found in environment variables")
    
    # Initialize OpenAI client
    client = OpenAI(api_key=api_key)
    
    # Determine file type
    file_type = get_file_type(file_path)
    
    # Extract content based on file type
    if file_type == 'pdf':
        contents = extract_pdf_content(file_path, client)
    else:  # presentation
        contents = extract_ppt_content(file_path, client)
    
    return contents

def extract_info_from_upload(input_file, output_path):
    try:
        # Process document
        contents = process_document(input_file)
        
        # Save to file
        save_to_file(contents, output_path)
        
        print(f"Successfully processed document and saved content to {output_path}")
        
    except Exception as e:
        print(f"Error processing document: {str(e)}")

# # Configuration
# input_file = "DeepSeek_R1.pdf"  # Replace with your file path (PDF or PPT/PPTX)
# output_path = input_file.split(".")[0] + ".txt"

# extract_info_from_upload(input_file, output_path)